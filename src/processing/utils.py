import chardet
from src.processing.pde_ple import PLE
from src.processing.document import Document
import pymupdf
import os
import subprocess
from pptx import Presentation
import pyxlsb
import ezodf
import csv
import pandas as pd
import re


def detect_encoding(file_path):
    with open(file_path, "rb") as f:
        result = chardet.detect(f.read())
    return result["encoding"]


def get_all_file_paths(root_path):
    useful_extensions = {
        ".pdf",
        ".doc",
        ".docx",
        ".txt",
        ".rtf",
        ".odt",
        ".ppt",
        ".pptx",
        ".pps",
        ".ppsx",
        ".odp",
        ".xls",
        ".xlsx",
        ".xlsb",
        ".ods",
        ".csv",
    }

    seen_paths = set()
    all_file_paths = set()
    file_count = 0

    stack = [(root_path, 0)]

    while stack:
        current_path, depth = stack.pop()
        if depth > 40 or current_path in seen_paths:
            continue
        seen_paths.add(current_path)

        try:
            df = PLE.s3.ls(current_path)
            paths = set(df["path"].to_list())  # deduplication at this level
        except Exception as e:
            print(f"Error listing {current_path}: {e}")
            continue

        for path in paths:
            file_name = path.split("/")[-1]
            if "." in file_name:
                ext = "." + file_name.split(".")[-1].lower()
                if ext in useful_extensions:
                    if path not in all_file_paths:
                        all_file_paths.add(path)
                        file_count += 1
                        if file_count % 2000 == 0:
                            print(f"Reached {file_count} files")
                continue
            else:
                # Assume it's a directory
                stack.append((path, depth + 1))

    return ["s3://" + p for p in all_file_paths]

def read_file_with_fallback(file_path, encodings):
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        raise UnicodeDecodeError(
            f"Could not decode the file with any of the tried encodings: {encodings}"
        )

def process_txt_file(file_path, encodings=["utf-8", "latin1"]):

    try:
        # Read the file with encoding fallback
        content = read_file_with_fallback(file_path, encodings)

        # Apply the remove_extra_spaces function
        processed_content = clean_text(content)

        # Save the processed content back to the file
        with open(file_path, "w", encoding="utf-8") as file:
            file.write(processed_content)

    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
        raise


def convert_to_docx(doc: Document):
    try:
        tempdir=os.path.dirname(doc.local_path)
        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "docx",
                doc.local_path,
                "--outdir",
                tempdir,
                "-env:UserInstallation=file:///tmp/LibreOffice_Conversion_${USER}",
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            check=True,
        )
    except Exception as e:
        print(f"Error converting to docx from {doc.s3_path}: {e}")


def convert_to_xlsx(doc: Document):
    try:
        tempdir = os.path.dirname(doc.local_path)
        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "xlsx",
                doc.local_path,
                "--outdir",
                tempdir,
                "-env:UserInstallation=file:///tmp/LibreOffice_Conversion_${USER}",
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            check=True,
        )
    except Exception as e:
        print(f"Error converting to xlsx from {doc.s3_path}: {e}")


def convert_to_pptx(doc: Document):
    try:
        tempdir = os.path.dirname(doc.local_path)
        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pptx",
                doc.local_path,
                "--outdir",
                tempdir,
                "-env:UserInstallation=file:///tmp/LibreOffice_Conversion_${USER}",
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            check=True,
        )
    except Exception as e:
        print(f"Error converting to pptx from {doc.s3_path}: {e}")



def extract_pptx(doc: Document):
    try:
        if not doc.extension=='.pptx':
            convert_to_pptx(doc)
        out = open(doc.local_txt_path, "w")

        prs = Presentation(doc.local_pptx_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(clean_text(shape.text))
        out.write("\n".join(text_runs))
        out.close()
    except Exception as e:
        print(f"Error extracting from {doc.s3_path}: {e}")

def extract_pdf_docx(doc: Document):
    path = doc.local_docx_path
    if doc.extension =='.pdf':
        path=doc.local_path
    
    try:
        document = pymupdf.open(
            path
        )  # open a document
        
        out = open(doc.local_txt_path, "wb")  # create a text output
        for page in document:  # iterate the document pages
            text = page.get_text()  # get plain text (is in UTF-8)
            text = clean_text(text)
            text=text.encode("utf8")
            out.write(text)  # write text of page

        out.close()

    except Exception as e:
        print(f"Error extracting from {doc.s3_path}: {e}")


def extract_doc(doc: Document):
    convert_to_docx(doc)
    extract_pdf_docx(doc)
    

def extract_xlsx(doc: Document):
    file_path = doc.local_path
    text_output = ""

    try:
        if doc.extension == ".csv":
            with open(file_path, newline="") as csvfile:
                reader = csv.reader(csvfile)
                for row in reader:
                    text_output += "\t".join(row) + "\n"

        elif doc.extension == ".xlsb":
            # Use pyxlsb only for .xlsb files
            with open(file_path, "rb") as f:
                df = pd.read_excel(f, engine="pyxlsb")
                text_output = df.to_string(index=False)

        else:
            if doc.extension in {".ods",".xls"}:
                convert_to_xlsx(doc)
                file_path=doc.local_xlsx_path
            xl = pd.ExcelFile(file_path)
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                text_output += f"Sheet: {sheet_name}\n"
                text_output += df.to_string(index=False, na_rep="") + "\n"

        # Write the collected text to a .txt file
        with open(doc.local_txt_path, "w", encoding="utf-8") as f:
            text_output = clean_text(text_output)
            f.write(text_output)

    except Exception as e:
        print(f"Error extracting from {doc.s3_path}: {e}")



def clean_text(text):
    # Replace non-breaking spaces with regular spaces
    text = text.replace("\xa0", " ")

    # Replace multiple newlines with a single newline
    text = re.sub(r"\n\s*\n", "\n", text)

    # Remove leading and trailing whitespace from each line
    lines = [line.strip() for line in text.splitlines()]

    # Join the lines back together with single newlines
    text = "\n".join(lines)

    # Use the existing function to remove extra spaces
    text = " ".join(text.split())

    return text