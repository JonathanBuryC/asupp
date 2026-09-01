from dataclasses import dataclass, field
from pathlib import Path
from src.constants.constants import USEFUL_EXTENSIONS,TEMP_DIR
from src.processing.pde_ple import PDE
import os
import subprocess
from pptx import Presentation
import pymupdf
import chardet
import re
import pandas as pd
import csv


@dataclass
class Document:
    s3_path: str
    s3_txt_path : str = ""
    local_path : str =""
    local_txt_path : str = ""
    local_docx_path: str = ""
    local_pptx_path: str = ""
    local_xlsx_path: str = ""
    extension: str = field(init=False)
    text: str = ""

    def __post_init__(self):
        self.extension = Path(self.s3_path).suffix.lower()
        self.filename = Path(self.s3_path).name.lower()
        root_s3, _ = os.path.splitext(self.s3_path)
        self.s3_txt_path = root_s3 + ".txt"
        self.local_path = TEMP_DIR + "/" + self.s3_path.split("/")[-1].lower()
        root_local, _ = os.path.splitext(self.local_path)
        root_local=root_local.lower()
        if self.local_path[-4::].lower()==".txt":
            self.local_txt_path=self.local_path
        else:
            self.local_txt_path = self.local_path + ".txt"   #pcq ya plusieurs doc avec mm noms et <> extensions donc l'extension dans ce cas doit rester dans le nom 
        self.local_docx_path = root_local + ".docx"
        self.raw_filename=Path(self.local_txt_path).name.lower()
        self.local_pptx_path = root_local + ".pptx"
        self.local_xlsx_path = root_local + ".xlsx"

    
    def download(self,dir=TEMP_DIR):
        if dir == TEMP_DIR and not os.path.exists(dir):
            os.makedirs(dir)
        PDE.s3.download(
            path=self.s3_path,
            local_file=dir + "/" + self.s3_path.split("/")[-1].lower(),
        )
        #print("successfully downloaded to :", self.local_path.lower())

    def upload_txt(self,s3_dir=None):
        if not s3_dir:
            s3_path=self.s3_txt_path
        else:
            s3_path=s3_dir+'/'+ self.local_txt_path.split('/')[-1]
        PDE.s3.upload(local_file=self.local_txt_path,path=s3_path)
        #print("successfully uploaded to :", s3_path)

    
     #def upload_txt(self,dir)
    def delete(self):
        if os.path.exists(self.local_path):
            os.remove(self.local_path)
        if os.path.exists(self.local_txt_path):
            os.remove(self.local_txt_path)
        if os.path.exists(self.local_docx_path):
            os.remove(self.local_docx_path)
        if os.path.exists(self.local_xlsx_path):
            os.remove(self.local_xlsx_path)
        if os.path.exists(self.local_pptx_path):
            os.remove(self.local_pptx_path)
        #print("successfully deleted :", self.local_path)

    def is_useful(self):
        return self.extension.lower() in USEFUL_EXTENSIONS
    
    def extract_text_by_extension(self) -> None:
        if self.extension in {".pdf", ".docx"}:
            self.extract_pdf_docx()
        elif self.extension in {".doc", ".odt", ".rtf"}:
            self.extract_doc()
        elif self.extension in {".ppt", ".pptx", ".pps", ".ppsx", ".odp"}:
            self.extract_pptx()
        elif self.extension in {".xls", ".xlsx", ".xlsm", ".ods", ".xlsb", ".csv"}:
            self.extract_xlsx()
        elif self.extension in {".txt"}:
            pass
        else:
            raise ValueError(f"Unsupported extension: {self.extension}")
        
        try:
            with open(self.local_txt_path, "r", encoding="utf-8") as f:
                self.text = f.read()
        except UnicodeDecodeError:
            with open(self.local_txt_path, "r", encoding="latin-1") as f:
                self.text = f.read()


    def convert_to_docx(self):
        try:
            tempdir = os.path.dirname(self.local_path)
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "docx",
                    self.local_path,
                    "--outdir",
                    tempdir,
                    "-env:UserInstallation=file:///tmp/LibreOffice_Conversion_${USER}",
                ],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                check=True,
            )
        except Exception as e:
            print(f"Error converting to docx from {self.s3_path}: {e}")

    def convert_to_xlsx(self):
        try:
            tempdir = os.path.dirname(self.local_path)
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "xlsx",
                    self.local_path,
                    "--outdir",
                    tempdir,
                    "-env:UserInstallation=file:///tmp/LibreOffice_Conversion_${USER}",
                ],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                check=True,
            )
        except Exception as e:
            print(f"Error converting to xlsx from {self.s3_path}: {e}")

    def convert_to_pptx(self):
        try:
            tempdir = os.path.dirname(self.local_path)
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pptx",
                    self.local_path,
                    "--outdir",
                    tempdir,
                    "-env:UserInstallation=file:///tmp/LibreOffice_Conversion_${USER}",
                ],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                check=True,
            )
        except Exception as e:
            print(f"Error converting to pptx from {self.s3_path}: {e}")

    def extract_pptx(self):
        try:
            if self.extension != ".pptx":
                self.convert_to_pptx()
            out = open(self.local_txt_path, "w")

            prs = Presentation(self.local_pptx_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(self.clean_text(shape.text))
            out.write("\n".join(text_runs))
            out.close()
        except Exception as e:
            print(f"Error extracting from {self.s3_path}: {e}")

    def extract_pdf_docx(self):
        path = self.local_docx_path if not self.extension == ".pdf" else self.local_path

        try:
            document = pymupdf.open(path)
            with open(self.local_txt_path, "wb") as out:
                for page in document:
                    text = page.get_text()
                    text = self.clean_text(text).encode("utf8")
                    out.write(text)
        except Exception as e:
            print(f"Error extracting from {self.s3_path}: {e}")

    def extract_doc(self):
        self.convert_to_docx()
        self.extract_pdf_docx()

    def extract_xlsx(self):
        file_path = self.local_path
        text_output = ""

        try:
            if self.extension == ".csv":
                with open(file_path, newline="") as csvfile:
                    reader = csv.reader(csvfile)
                    for row in reader:
                        text_output += "\t".join(row) + "\n"

            elif self.extension == ".xlsb":
                with open(file_path, "rb") as f:
                    df = pd.read_excel(f, engine="pyxlsb")
                    text_output = df.to_string(index=False)

            else:
                if self.extension in {".ods", ".xls"}:
                    self.convert_to_xlsx()
                    file_path = self.local_xlsx_path
                xl = pd.ExcelFile(file_path)
                for sheet_name in xl.sheet_names:
                    df = xl.parse(sheet_name)
                    text_output += f"Sheet: {sheet_name}\n"
                    text_output += df.to_string(index=False, na_rep="") + "\n"

            with open(self.local_txt_path, "w", encoding="utf-8") as f:
                text_output = self.clean_text(text_output)
                f.write(text_output)

        except Exception as e:
            print(f"Error extracting from {self.s3_path}: {e}")

    def clean_text(self, text: str) -> str:
        text = text.replace("\xa0", " ")

        # Replace multiple newlines with a single newline
        text = re.sub(r"\n\s*\n", "\n", text)

        # Remove leading and trailing whitespace from each line
        lines = [line.strip() for line in text.splitlines()]

        # Join the lines back together with single newlines
        text = "\n".join(lines)

        # Use the existing function to remove extra spaces
        text = " ".join(text.split())

        return text
        
    
